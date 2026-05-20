import base64
import json
import os
import subprocess
import tempfile
from io import BytesIO
from pathlib import Path
import re
from uuid import uuid4

from langchain_core.documents import Document
from langchain_huggingface import HuggingFaceEmbeddings
from huggingface_hub import InferenceClient
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, PayloadSchemaType, PointStruct, VectorParams

from ingestion import (
    EMBEDDING_MODEL,
    EMBEDDING_VECTOR_SIZE,
    PPTX_MIME,
    QDRANT_API_KEY,
    QDRANT_COLLECTION_NAME,
    QDRANT_URL,
    build_drive_service,
    download_drive_file,
    infer_excel_match_column,
    load_drive_file_metadata,
    load_reference_excel_rows,
    map_drive_ppts_to_excel_rows_exact,
    normalize_match_key,
)


DEFAULT_QDRANT_COLLECTION_NAME = "predikly_t7"
MAX_USECASE_CONTENT_CHARS = 7000
CUSTOMER_MANIFEST_PATH = Path("data/customer_manifest.json")

VISION_MODEL = "google/gemma-4-31B-it"

VISION_PROMPT = """You are analyzing a PowerPoint slide image from a Predikly customer case-study deck.

Summarize the visible workflow, diagram, architecture, or flowchart.

Return JSON only with this structure:

{
  "workflow_image_summary": "",
  "solution_proposed": "",
  "tools_used": [],
  "benefits": []
}

Rules:
- Describe only what is visible or strongly supported by the slide.
- Do not invent tools, technologies, benefits, or implementation details.
- If a field is not visible or cannot be inferred safely, return an empty string or empty list.
- Keep workflow_image_summary clear and concise.
- Mention major entities, process steps, systems, arrows, decisions, outputs, and dashboards if visible.
"""


TOOLS_LABELS = (
    "tools used",
    "technology used",
    "technologies used",
    "applications used",
    "systems used",
    "platforms used",
    "technology stack",
    "tech stack",
    "tools",
)

PROBLEM_LABELS = (
    "problem statement",
    "business problem",
    "challenge",
    "challenges",
    "pain points",
    "current process challenges",
)

SOLUTION_LABELS = (
    "solution summary",
    "solution",
    "proposed solution",
    "solution approach",
    "automation solution",
    "target solution",
    "rpa solution",
)

BENEFIT_LABELS = (
    "benefits",
    "business benefits",
    "outcomes",
    "impact",
    "results",
    "outputs",
    "business roi",
    "roi",
    "business value",
    "value delivered",
)

USE_CASE_CATEGORY_COLUMNS = (
    "Use Case Category",
    "UseCase Category",
    "Category",
)


def normalize_customer_name(value: str) -> str:
    value = value.lower().strip()
    value = re.sub(r"[^a-z0-9\s]", " ", value)
    return re.sub(r"\s+", " ", value).strip()


def clean_metadata_value(value):
    if value is None:
        return ""

    if isinstance(value, (str, int, float, bool)):
        return value

    return str(value)


def clean_excel_row(row: dict) -> dict:
    return {
        str(key): clean_metadata_value(value)
        for key, value in row.items()
    }


def excel_value(row: dict, *column_names: str) -> str:
    normalized_row = {
        str(key).strip().lower(): clean_metadata_value(value)
        for key, value in row.items()
    }

    for column_name in column_names:
        value = normalized_row.get(column_name.strip().lower())

        if value not in (None, ""):
            return str(value).strip()

    return ""


def one_line(value: str) -> str:
    return re.sub(r"\s+", " ", value).strip()


def line_starts_with_label(line: str, labels: tuple[str, ...]) -> bool:
    normalized_line = line.lower().strip()
    return any(normalized_line.startswith(label) for label in labels)


def is_known_section_label(line: str) -> bool:
    return any(
        line_starts_with_label(line, labels)
        for labels in (
            TOOLS_LABELS,
            PROBLEM_LABELS,
            SOLUTION_LABELS,
            BENEFIT_LABELS,
        )
    )


def extract_labeled_value(
    lines: list[str],
    labels: tuple[str, ...],
    max_lines: int = 8,
) -> str:
    for index, line in enumerate(lines):
        stripped_line = line.strip()

        if not line_starts_with_label(stripped_line, labels):
            continue

        if ":" in stripped_line:
            _, value = stripped_line.split(":", 1)

            if value.strip():
                inline_value = one_line(value)

                captured_after_inline = []

                for next_line in lines[index + 1:index + 1 + max_lines]:
                    stripped_next_line = next_line.strip()

                    if not stripped_next_line:
                        continue

                    if is_known_section_label(stripped_next_line):
                        break

                    captured_after_inline.append(stripped_next_line)

                if captured_after_inline:
                    return one_line(
                        " ".join([inline_value, *captured_after_inline])
                    )

                return inline_value

        captured_lines = []

        for next_line in lines[index + 1:index + 1 + max_lines]:
            stripped_next_line = next_line.strip()

            if not stripped_next_line:
                continue

            if is_known_section_label(stripped_next_line):
                break

            captured_lines.append(stripped_next_line)

        return one_line(" ".join(captured_lines))

    return ""


def split_labeled_list(value: str) -> list[str]:
    if not value:
        return []

    value = value.replace("•", "\n")
    value = value.replace("–", "-")
    value = value.replace("—", "-")
    value = re.sub(r"\s{2,}", "\n", value)

    items = re.split(
        r"\n|;|(?:\s+-\s+)|"
        r"(?=\bTime-Saving\b)|"
        r"(?=\bTime Saving\b)|"
        r"(?=\bIncreased Efficiency\b)|"
        r"(?=\bCost Saving\b)|"
        r"(?=\bCost Reduction\b)|"
        r"(?=\$ saving\b)|"
        r"(?=\bSavings\b)|"
        r"(?=\bReduced\b)|"
        r"(?=\bImproved\b)|"
        r"(?=\bFaster\b)|"
        r"(?=\bBetter\b)",
        value,
        flags=re.IGNORECASE,
    )

    cleaned_items = []

    for item in items:
        cleaned = one_line(item)
        cleaned = cleaned.strip(":- ")

        if cleaned and cleaned not in cleaned_items:
            cleaned_items.append(cleaned)

    return cleaned_items


def extract_ppt_metadata_fields(slide_texts: list[tuple[int, str]]) -> dict:
    all_lines = []

    for _, slide_text in slide_texts:
        all_lines.extend(slide_text.splitlines())

    return {
        "tools_used": split_labeled_list(
            extract_labeled_value(all_lines, TOOLS_LABELS, max_lines=6)
        ),
        "problem_statement": extract_labeled_value(
            all_lines,
            PROBLEM_LABELS,
            max_lines=10,
        ),
        "solution_summary": extract_labeled_value(
            all_lines,
            SOLUTION_LABELS,
            max_lines=10,
        ),
        "benefits": split_labeled_list(
            extract_labeled_value(all_lines, BENEFIT_LABELS, max_lines=10)
        ),
    }


def extract_slide_title(slide, slide_text: str) -> str:
    if slide.shapes.title and getattr(slide.shapes.title, "text", "").strip():
        return one_line(slide.shapes.title.text)

    for line in slide_text.splitlines():
        title = one_line(line)

        if title:
            return title

    return ""


def extract_ppt_slide_text(slide) -> str:
    slide_lines = []

    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            slide_lines.append(shape.text.strip())

    return "\n".join(slide_lines).strip()


def extract_pptx_slide_infos(content: bytes) -> list[dict]:
    presentation = Presentation(BytesIO(content))
    slide_infos = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = extract_ppt_slide_text(slide)

        slide_infos.append(
            {
                "slide_number": slide_number,
                "slide_title": extract_slide_title(slide, slide_text),
                "visible_text": slide_text,
                "has_diagram": detect_diagram_or_image(slide),
                "content_type": detect_content_type(slide, slide_text),
            }
        )

    return slide_infos


def extract_pptx_slide_texts(content: bytes) -> list[tuple[int, str]]:
    return [
        (slide_info["slide_number"], slide_info["visible_text"])
        for slide_info in extract_pptx_slide_infos(content)
        if slide_info["visible_text"]
    ]


def detect_diagram_or_image(slide) -> bool:
    diagram_shape_count = 0
    arrow_or_connector_count = 0
    text_box_count = 0

    for shape in slide.shapes:
        if shape.shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.GROUP,
            MSO_SHAPE_TYPE.DIAGRAM,
            MSO_SHAPE_TYPE.CHART,
        }:
            return True

        if getattr(shape, "has_table", False):
            return True

        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            diagram_shape_count += 1

            if hasattr(shape, "text") and shape.text.strip():
                text_box_count += 1
                text = shape.text.strip().lower()

                if any(
                    term in text
                    for term in (
                        "workflow",
                        "process",
                        "architecture",
                        "flow",
                        "bot",
                        "automation",
                        "dashboard",
                        "system",
                        "input",
                        "output",
                    )
                ):
                    return True

        shape_name = getattr(shape, "name", "").lower()

        if "arrow" in shape_name or "connector" in shape_name:
            arrow_or_connector_count += 1

    if diagram_shape_count >= 4 and text_box_count >= 2:
        return True

    if arrow_or_connector_count >= 1 and text_box_count >= 2:
        return True

    return False


def detect_content_type(slide, slide_text: str) -> str:
    has_picture = False
    has_chart = False
    has_table = False
    has_diagram = False

    for shape in slide.shapes:
        has_picture = has_picture or shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        has_chart = has_chart or shape.shape_type == MSO_SHAPE_TYPE.CHART
        has_table = has_table or getattr(shape, "has_table", False)
        has_diagram = has_diagram or shape.shape_type in {
            MSO_SHAPE_TYPE.DIAGRAM,
            MSO_SHAPE_TYPE.GROUP,
        }

    lowered = slide_text.lower()

    if has_chart:
        return "chart"

    if has_table:
        return "table"

    if has_diagram or any(
        term in lowered
        for term in (
            "architecture",
            "system architecture",
            "technical architecture",
            "solution architecture",
        )
    ):
        return "architecture_diagram"

    if has_picture or any(
        term in lowered
        for term in (
            "workflow",
            "flowchart",
            "process flow",
            "process",
            "bot",
            "automation",
            "rpa",
        )
    ):
        return "workflow_diagram"

    return "text"


def render_slide_image(ppt_content: bytes, slide_number: int) -> bytes | None:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        ppt_path = temp_path / "deck.pptx"
        output_dir = temp_path / "rendered"
        output_dir.mkdir()
        ppt_path.write_bytes(ppt_content)

        try:
            subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "png",
                    "--outdir",
                    str(output_dir),
                    str(ppt_path),
                ],
                check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.DEVNULL,
            )
        except FileNotFoundError:
            print("LibreOffice/soffice not found. Slide rendering skipped.")
            return None
        except subprocess.CalledProcessError as error:
            print(f"Slide rendering failed: {error}")
            return None

        rendered_images = sorted(output_dir.glob("*.png"))

        if not rendered_images:
            print("No rendered slide images found.")
            return None

        if slide_number > len(rendered_images):
            print(
                f"Rendered image count is {len(rendered_images)}, "
                f"but requested slide {slide_number}."
            )
            return None

        return rendered_images[slide_number - 1].read_bytes()


def extract_first_slide_image(content: bytes, slide_number: int) -> bytes | None:
    presentation = Presentation(BytesIO(content))

    if slide_number < 1 or slide_number > len(presentation.slides):
        return None

    slide = presentation.slides[slide_number - 1]

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return shape.image.blob

    return None


def render_or_extract_slide_image(ppt_content: bytes, slide_number: int) -> bytes | None:
    return render_slide_image(ppt_content, slide_number) or extract_first_slide_image(
        ppt_content,
        slide_number,
    )


def safe_parse_json_from_llm(content: str) -> dict:
    if not content:
        return {}

    content = content.strip()

    try:
        return json.loads(content)
    except json.JSONDecodeError:
        pass

    match = re.search(r"\{.*\}", content, flags=re.DOTALL)

    if not match:
        return {}

    try:
        return json.loads(match.group(0))
    except json.JSONDecodeError:
        return {}


def call_vision_llm(image_bytes: bytes) -> dict:
    empty_result = {
        "workflow_image_summary": "",
        "solution_proposed": "",
        "tools_used": [],
        "benefits": [],
    }

    token = os.getenv("HF_HUB_TOKEN") or os.getenv("HUGGINGFACEHUB_API_TOKEN")

    if not token:
        print("HF token missing. Vision LLM skipped.")
        return empty_result

    image_data = base64.b64encode(image_bytes).decode("utf-8")
    client = InferenceClient(model=VISION_MODEL, token=token)

    try:
        response = client.chat_completion(
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": VISION_PROMPT},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_data}",
                            },
                        },
                    ],
                }
            ],
            max_tokens=700,
        )

        content = response.choices[0].message.content
        parsed = safe_parse_json_from_llm(content)

        if not parsed:
            print(f"Vision LLM returned non-JSON content: {content}")
            return empty_result

    except Exception as error:
        print(f"Vision LLM failed: {error}")
        return empty_result

    tools_used = parsed.get("tools_used", [])
    benefits = parsed.get("benefits", [])

    return {
        "workflow_image_summary": clean_metadata_value(
            parsed.get("workflow_image_summary", "")
        ),
        "solution_proposed": clean_metadata_value(
            parsed.get("solution_proposed", "")
        ),
        "tools_used": tools_used if isinstance(tools_used, list) else [],
        "benefits": benefits if isinstance(benefits, list) else [],
    }


def tokenize_for_match(value: str) -> set[str]:
    stop_words = {
        "a",
        "an",
        "and",
        "for",
        "in",
        "of",
        "on",
        "or",
        "the",
        "to",
        "use",
        "case",
        "process",
        "automation",
    }

    return {
        token
        for token in re.findall(r"[a-z0-9]+", value.lower())
        if len(token) > 2 and token not in stop_words
    }


def select_usecase_slide_texts(
    slide_texts: list[tuple[int, str]],
    *,
    usecase_name: str,
    total_usecases_for_ppt: int,
) -> list[tuple[int, str]]:
    if total_usecases_for_ppt <= 1:
        return slide_texts

    tokens = tokenize_for_match(usecase_name)

    if not tokens:
        return slide_texts

    scored_slides = []

    for slide_number, slide_text in slide_texts:
        slide_tokens = tokenize_for_match(slide_text)
        score = len(tokens.intersection(slide_tokens))

        if score:
            scored_slides.append((score, slide_number, slide_text))

    if not scored_slides:
        return slide_texts

    scored_slides.sort(key=lambda item: (-item[0], item[1]))

    selected = [
        (slide_number, slide_text)
        for _, slide_number, slide_text in scored_slides[:4]
    ]

    selected.sort(key=lambda item: item[0])
    return selected


def select_usecase_slide_infos(
    slide_infos: list[dict],
    *,
    usecase_name: str,
    total_usecases_for_ppt: int,
) -> list[dict]:
    if total_usecases_for_ppt <= 1:
        return slide_infos

    selected_slide_numbers = {
        slide_number
        for slide_number, _ in select_usecase_slide_texts(
            [
                (slide_info["slide_number"], slide_info["visible_text"])
                for slide_info in slide_infos
            ],
            usecase_name=usecase_name,
            total_usecases_for_ppt=total_usecases_for_ppt,
        )
    }

    return [
        slide_info
        for slide_info in slide_infos
        if slide_info["slide_number"] in selected_slide_numbers
    ] or slide_infos


def build_usecase_page_content(
    *,
    excel_row: dict,
    metadata: dict,
    selected_slide_texts: list[tuple[int, str]],
) -> str:
    parts = [
        f"Customer Name: {metadata['customer_name']}",
        f"Customer Domain: {metadata['customer_domain']}",
        f"Use Case Name: {metadata['usecase_name']}",
        f"Tools Used: {metadata['tools_used']}",
        f"Problem Statement: {metadata['problem_statement']}",
        f"Solution Summary: {metadata['solution_summary']}",
        f"Use Case Link: {clean_metadata_value(excel_row.get('Use Case Link', ''))}",
        "",
        "Relevant PPT Content:",
    ]

    for slide_number, slide_text in selected_slide_texts:
        parts.append(f"[Slide {slide_number}]")
        parts.append(slide_text)

    return "\n".join(parts).strip()[:MAX_USECASE_CONTENT_CHARS]


def build_content_for_embedding(metadata: dict, visible_text: str = "") -> str:
    return "\n".join(
        [
            f"Company: {metadata.get('company_name', '')}",
            f"Customer Domain: {metadata.get('customer_domain', '')}",
            f"Use Case: {metadata.get('use_case_name', '')}",
            f"Use Case Category: {metadata.get('use_case_category', '')}",
            f"Solution Proposed: {metadata.get('solution_proposed', '')}",
            f"Workflow/Image Summary: {metadata.get('workflow_image_summary', '')}",
            f"Tools Used: {', '.join(metadata.get('tools_used', []))}",
            f"Benefits: {', '.join(metadata.get('benefits', []))}",
            "",
            "Visible Slide Text:",
            visible_text,
        ]
    ).strip()[:MAX_USECASE_CONTENT_CHARS]


def extract_slide_fields_from_text(slide_text: str) -> dict:
    lines = slide_text.splitlines()

    return {
        "solution_proposed": extract_labeled_value(
            lines,
            SOLUTION_LABELS,
            max_lines=8,
        ),
        "tools_used": split_labeled_list(
            extract_labeled_value(
                lines,
                TOOLS_LABELS,
                max_lines=5,
            )
        ),
        "benefits": split_labeled_list(
            extract_labeled_value(
                lines,
                BENEFIT_LABELS,
                max_lines=8,
            )
        ),
    }


def merge_list_fields(primary: list[str], fallback: list[str]) -> list[str]:
    merged = []

    for item in [*primary, *fallback]:
        cleaned = one_line(str(item))

        if cleaned and cleaned not in merged:
            merged.append(cleaned)

    return merged


def infer_tools_from_text(slide_text: str) -> list[str]:
    lowered = slide_text.lower()
    tools = []

    tool_patterns = {
        "RPA Bot": ["rpa", "bot", "bots"],
        "HL7 Data": ["hl7"],
        "Client In-house Application": [
            "in house application",
            "in-house application",
            "client application",
        ],
        "Insurance Provider Website": [
            "insurance provider",
            "provider website",
            "website",
        ],
        "Workflow Automation": [
            "automated",
            "automation",
            "automate",
            "automated process",
        ],
        "Dashboard": ["dashboard", "reporting"],
        "OCR": ["ocr"],
        "API Integration": ["api"],
        "CRM": ["crm"],
        "Email Automation": ["email"],
        "Document Processing": ["document processing", "file extraction", "files extraction"],
    }

    for tool, patterns in tool_patterns.items():
        if any(pattern in lowered for pattern in patterns):
            tools.append(tool)

    return tools


def infer_benefits_from_text(slide_text: str) -> list[str]:
    lines = [one_line(line) for line in slide_text.splitlines() if one_line(line)]
    benefits = []

    benefit_keywords = (
        "saving",
        "cost reduction",
        "profitability",
        "time-saving",
        "time saving",
        "increased efficiency",
        "efficiency",
        "reduced",
        "reduction",
        "faster",
        "accurate",
        "accuracy",
        "productivity",
        "roi",
    )

    for line in lines:
        lowered = line.lower()

        if any(keyword in lowered for keyword in benefit_keywords):
            cleaned = line.strip(":- ")

            if cleaned and cleaned not in benefits:
                benefits.append(cleaned)

    return benefits


def build_vector_metadata_payload(
    *,
    vector_point_id: str,
    excel_row: dict,
    drive_file: dict,
    match_key: str,
    slide_info: dict,
    ppt_fields: dict,
    vision_fields: dict,
) -> dict:
    company_name = (
        excel_value(excel_row, "Customer Name", "Company Name")
        or clean_metadata_value(drive_file.get("name", ""))
    )

    customer_domain = excel_value(excel_row, "Customer Domain")
    use_case_name = excel_value(excel_row, "Use Case Name")
    use_case_category = excel_value(excel_row, *USE_CASE_CATEGORY_COLUMNS)

    visible_text = slide_info.get("visible_text", "")
    text_fields = extract_slide_fields_from_text(visible_text)

    solution_proposed = (
        text_fields["solution_proposed"]
        or vision_fields.get("solution_proposed", "")
        or ppt_fields.get("solution_summary", "")
    )

    tools_from_text = text_fields["tools_used"] or ppt_fields.get("tools_used", [])
    tools_from_vision = vision_fields.get("tools_used", [])
    tools_inferred = infer_tools_from_text(visible_text)

    tools_used = merge_list_fields(
        merge_list_fields(tools_from_text, tools_from_vision),
        tools_inferred,
    )

    benefits_from_text = text_fields["benefits"] or ppt_fields.get("benefits", [])
    benefits_from_vision = vision_fields.get("benefits", [])
    benefits_inferred = infer_benefits_from_text(visible_text)

    benefits = merge_list_fields(
        merge_list_fields(benefits_from_text, benefits_from_vision),
        benefits_inferred,
    )

    workflow_image_summary = clean_metadata_value(
        vision_fields.get("workflow_image_summary", "")
    )

    content_type = slide_info.get("content_type", "text")

    if workflow_image_summary:
        chunk_type = "diagram_summary"
    elif content_type in ("workflow_diagram", "architecture_diagram"):
        chunk_type = "slide_summary"
    else:
        chunk_type = "slide_text"

    metadata = {
        "vector_point_id": vector_point_id,
        "document_id": clean_metadata_value(drive_file.get("id", "")),
        "ppt_name": clean_metadata_value(drive_file.get("name", "")),
        "slide_number": slide_info.get("slide_number"),
        "slide_title": clean_metadata_value(slide_info.get("slide_title", "")),
        "company_name": company_name,
        "customer_name_normalized": normalize_customer_name(company_name),
        "customer_domain": customer_domain,
        "use_case_name": use_case_name,
        "use_case_category": use_case_category,
        "solution_proposed": solution_proposed,
        "workflow_image_summary": workflow_image_summary,
        "tools_used": tools_used,
        "benefits": benefits,
        "chunk_type": chunk_type,
        "content_type": content_type,
        "source_type": "internal_ppt",
        "is_internal": True,
    }

    return metadata


def build_usecase_document(
    *,
    excel_row: dict,
    metadata: dict,
    selected_slide_texts: list[tuple[int, str]],
) -> Document:
    return Document(
        page_content=build_usecase_page_content(
            excel_row=excel_row,
            metadata=metadata,
            selected_slide_texts=selected_slide_texts,
        ),
        metadata=metadata,
    )


def build_usecase_metadata(
    *,
    excel_row: dict,
    drive_file: dict,
    match_key: str,
    ppt_fields: dict,
) -> dict:
    return {
        "customer_domain": excel_value(excel_row, "Customer Domain"),
        "usecase_name": excel_value(excel_row, "Use Case Name"),
        "customer_name": excel_value(excel_row, "Customer Name", "Company Name"),
        "drive_id": drive_file["id"],
        "tools_used": ppt_fields["tools_used"],
        "problem_statement": ppt_fields["problem_statement"],
        "solution_summary": ppt_fields["solution_summary"],
        "benefits": ppt_fields["benefits"],
        "match_key": match_key,
    }


def build_documents_for_qdrant() -> list[Document]:
    service = build_drive_service()
    drive_files = load_drive_file_metadata(service)
    extra_excel_rows = []

    try:
        excel_rows = load_reference_excel_rows(service, drive_files)
        excel_match_column = infer_excel_match_column(excel_rows, drive_files)

        mapped_items, unmatched_ppts, extra_excel_rows = map_drive_ppts_to_excel_rows_exact(
            drive_files=drive_files,
            excel_rows=excel_rows,
            excel_match_column=excel_match_column,
        )
    except Exception as error:
        print(f"Excel metadata unavailable; continuing with PPT-only metadata. Reason: {error}")

        mapped_items = []
        unmatched_ppts = [
            {
                "drive_file": drive_file,
                "match_key": normalize_match_key(drive_file.get("name", "")),
            }
            for drive_file in drive_files
            if drive_file.get("mimeType") == PPTX_MIME
        ]

    documents = []

    ppt_items = [
        {
            "drive_file": item["drive_file"],
            "excel_rows": item["excel_rows"],
            "match_key": item["match_key"],
        }
        for item in mapped_items
    ]

    ppt_items.extend(
        {
            "drive_file": item["drive_file"],
            "excel_rows": [{}],
            "match_key": item["match_key"],
        }
        for item in unmatched_ppts
    )

    for item in ppt_items:
        drive_file = item["drive_file"]

        if drive_file.get("mimeType") != PPTX_MIME:
            continue

        print(f"\nProcessing PPT: {drive_file.get('name', '')}")

        ppt_content = download_drive_file(service, drive_file["id"])
        slide_infos = extract_pptx_slide_infos(ppt_content)

        slide_texts = [
            (slide_info["slide_number"], slide_info["visible_text"])
            for slide_info in slide_infos
            if slide_info["visible_text"]
        ]

        ppt_fields = extract_ppt_metadata_fields(slide_texts)

        for excel_row in item["excel_rows"]:
            base_metadata = build_usecase_metadata(
                excel_row=excel_row,
                drive_file=drive_file,
                match_key=item["match_key"],
                ppt_fields=ppt_fields,
            )

            selected_slide_infos = select_usecase_slide_infos(
                slide_infos,
                usecase_name=base_metadata["usecase_name"],
                total_usecases_for_ppt=len(item["excel_rows"]),
            )

            for slide_info in selected_slide_infos:
                vision_fields = {
                    "workflow_image_summary": "",
                    "solution_proposed": "",
                    "tools_used": [],
                    "benefits": [],
                }

                if slide_info["has_diagram"]:
                    print(
                        f"Diagram/image detected: "
                        f"{drive_file.get('name')} | slide {slide_info['slide_number']}"
                    )

                    image_bytes = render_or_extract_slide_image(
                        ppt_content,
                        slide_info["slide_number"],
                    )

                    if not image_bytes:
                        print(
                            f"No image rendered/extracted for slide "
                            f"{slide_info['slide_number']}"
                        )
                    else:
                        print(
                            f"Image bytes available for slide "
                            f"{slide_info['slide_number']}: {len(image_bytes)} bytes"
                        )

                        vision_fields = call_vision_llm(image_bytes)
                        print(f"Vision fields returned: {vision_fields}")

                vector_point_id = str(uuid4())

                metadata = build_vector_metadata_payload(
                    vector_point_id=vector_point_id,
                    excel_row=excel_row,
                    drive_file=drive_file,
                    match_key=item["match_key"],
                    slide_info=slide_info,
                    ppt_fields=ppt_fields,
                    vision_fields=vision_fields,
                )

                page_content = build_content_for_embedding(
                    metadata,
                    visible_text=slide_info.get("visible_text", ""),
                )

                documents.append(
                    Document(
                        page_content=page_content,
                        metadata=metadata,
                    )
                )

    print(f"\nMapped PPTs: {len(mapped_items)}")
    print(f"Extra Excel rows not used for this folder: {len(extra_excel_rows)}")
    print(f"Unmatched Drive PPTs ingested with PPT-only metadata: {len(unmatched_ppts)}")
    print(f"Slide vector documents for upload: {len(documents)}")

    return documents


def build_customer_manifest(documents: list[Document]) -> list[dict]:
    customers = {}

    for document in documents:
        metadata = document.metadata
        normalized_name = metadata.get("customer_name_normalized", "")

        if not normalized_name:
            continue

        customer = customers.setdefault(
            normalized_name,
            {
                "company_name": metadata.get("company_name", ""),
                "customer_name_normalized": normalized_name,
                "customer_domain": metadata.get("customer_domain", ""),
                "use_cases": {},
            },
        )
        use_case_name = metadata.get("use_case_name", "")

        if use_case_name:
            customer["use_cases"].setdefault(
                use_case_name,
                {
                    "use_case_name": use_case_name,
                    "customer_domain": metadata.get("customer_domain", ""),
                    "ppt_names": set(),
                    "slide_numbers": set(),
                },
            )
            use_case = customer["use_cases"][use_case_name]

            if metadata.get("ppt_name"):
                use_case["ppt_names"].add(metadata["ppt_name"])

            if metadata.get("slide_number") is not None:
                use_case["slide_numbers"].add(metadata["slide_number"])

    manifest = []

    for customer in customers.values():
        use_cases = []

        for use_case in customer["use_cases"].values():
            use_cases.append(
                {
                    **use_case,
                    "ppt_names": sorted(use_case["ppt_names"]),
                    "slide_numbers": sorted(use_case["slide_numbers"]),
                }
            )

        manifest.append(
            {
                **customer,
                "use_cases": sorted(
                    use_cases,
                    key=lambda item: item["use_case_name"].lower(),
                ),
                "use_case_count": len(use_cases),
            }
        )

    return sorted(manifest, key=lambda item: item["customer_name_normalized"])


def write_customer_manifest(documents: list[Document]) -> None:
    manifest = build_customer_manifest(documents)
    CUSTOMER_MANIFEST_PATH.parent.mkdir(parents=True, exist_ok=True)
    CUSTOMER_MANIFEST_PATH.write_text(
        json.dumps(manifest, indent=2),
        encoding="utf-8",
    )
    print(f"Wrote customer manifest: {CUSTOMER_MANIFEST_PATH} ({len(manifest)} customer(s))")


def chunk_documents(documents: list[Document]) -> list[Document]:
    return documents


def ensure_qdrant_collection(client: QdrantClient) -> None:
    collection_name = QDRANT_COLLECTION_NAME or DEFAULT_QDRANT_COLLECTION_NAME

    if client.collection_exists(collection_name):
        return

    client.create_collection(
        collection_name=collection_name,
        vectors_config=VectorParams(
            size=EMBEDDING_VECTOR_SIZE,
            distance=Distance.COSINE,
        ),
    )


def ensure_payload_indexes(client: QdrantClient, collection_name: str) -> None:
    indexes = {
        "customer_name_normalized": PayloadSchemaType.KEYWORD,
        "company_name": PayloadSchemaType.TEXT,
        "use_case_name": PayloadSchemaType.TEXT,
        "source_type": PayloadSchemaType.KEYWORD,
        "is_internal": PayloadSchemaType.BOOL,
        "content_type": PayloadSchemaType.KEYWORD,
    }

    for field_name, field_schema in indexes.items():
        try:
            client.create_payload_index(
                collection_name=collection_name,
                field_name=field_name,
                field_schema=field_schema,
            )
        except Exception as error:
            message = str(error).lower()

            if "already exists" not in message:
                print(f"Payload index skipped for {field_name}: {error}")


def upload_documents_to_qdrant(documents: list[Document]) -> None:
    if not QDRANT_URL or not QDRANT_API_KEY:
        raise ValueError("QDRANT_URL and QDRANT_API_KEY must be set in .env")

    if not documents:
        raise ValueError("No documents were created from the mapped PPTs.")

    collection_name = QDRANT_COLLECTION_NAME or DEFAULT_QDRANT_COLLECTION_NAME

    client = QdrantClient(
        url=QDRANT_URL,
        api_key=QDRANT_API_KEY,
    )

    ensure_qdrant_collection(client)
    ensure_payload_indexes(client, collection_name)

    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

    vectors = embeddings.embed_documents(
        [document.page_content for document in documents]
    )

    points = []

    for document, vector in zip(documents, vectors):
        payload = {
            **document.metadata,
            "page_content": document.page_content,
        }

        points.append(
            PointStruct(
                id=str(document.metadata.get("vector_point_id") or uuid4()),
                vector=vector,
                payload=payload,
            )
        )

    client.upsert(
        collection_name=collection_name,
        points=points,
    )

    print(
        f"Uploaded {len(documents)} vector point(s) "
        f"to Qdrant collection: {collection_name}"
    )


def main() -> None:
    documents = build_documents_for_qdrant()
    chunks = chunk_documents(documents)

    print(f"Vector points for upload: {len(chunks)}")
    write_customer_manifest(chunks)
    upload_documents_to_qdrant(chunks)


if __name__ == "__main__":
    main()
